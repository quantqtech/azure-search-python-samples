"""
From Ground to Air — Gent Machine
Built using deck_format.py for Aidan branding.
Bauhaus 93 lowercase titles in aidan blue, Century Gothic content, clean white.
"""

import sys, os
sys.stdout.reconfigure(encoding='utf-8')
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from deck_format import (
    AIDAN_BLUE, BLACK, WHITE, MAIN_BG, SECONDARY, DARK, COLORS,
    TABLE_ALT_ROW, FONT_TITLE, FONT_CONTENT,
    SLIDE_WIDTH, SLIDE_HEIGHT,
    LOGO_AIDAN, LOGO_AIDAN_ICON,
    set_slide_bg, add_title, add_subtitle, add_textbox, add_bullet_list,
    add_table, add_accent_line, add_slide_number, add_callout_box,
    add_icon
)

OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      "From Ground to Air - Gent Machine.pptx")

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank = prs.slide_layouts[6]
sn = 0

def content_slide(title_text, subtitle_text=None):
    """Create a standard content slide with title, accent line, icon, and slide number."""
    global sn
    sn += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title(s, title_text)
    add_accent_line(s)
    if subtitle_text:
        add_subtitle(s, subtitle_text)
    add_icon(s)
    add_slide_number(s, sn + 1)
    return s


# ============================================================
# 1. TITLE
# ============================================================
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, WHITE)

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "from ground to air"
p.font.size = Pt(54); p.font.color.rgb = AIDAN_BLUE
p.font.name = FONT_TITLE; p.alignment = PP_ALIGN.CENTER

add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "Gent Machine \u2014 Davenport Maintenance Assistant",
            font_size=24, font_color=BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
            "Phase 2: Managed AIR", font_size=20,
            font_color=AIDAN_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.6), Inches(2.3), Inches(0.025))
sh.fill.solid(); sh.fill.fore_color.rgb = AIDAN_BLUE; sh.line.fill.background()

add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
            "Aidan Systems  \u00b7  MAGNET  \u00b7  March 2026",
            font_size=16, font_color=BLACK, alignment=PP_ALIGN.CENTER)

if os.path.exists(LOGO_AIDAN):
    slide.shapes.add_picture(LOGO_AIDAN, Inches(5.0), Inches(5.8), height=Inches(1.0))


# ============================================================
# 2. AI — A NEW OPERATING ALTITUDE
# ============================================================
slide = content_slide("ai \u2014 a new operating altitude")

for i, (t, d) in enumerate([
    ("speed at scale", "Decisions that took days happen in seconds.\nWork runs autonomously."),
    ("intelligence that compounds", "Every interaction makes the system smarter.\nAI accelerates the more you use it."),
    ("reach without friction", "Real-time insight across operations \u2014\nwithout the coordination overhead."),
]):
    left = Inches(0.8) + i * Inches(4.1)
    add_callout_box(slide, left, Inches(1.8), Inches(3.6), Inches(0.7), t, font_size=16)
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.7), Inches(3.6), Inches(2.5))
    body.fill.solid(); body.fill.fore_color.rgb = TABLE_ALT_ROW
    body.line.color.rgb = COLORS["main"]; body.line.width = Pt(1)
    tf = body.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = d
    p.font.size = Pt(15); p.font.color.rgb = BLACK; p.font.name = FONT_CONTENT
    p.alignment = PP_ALIGN.CENTER

add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
            "This isn\u2019t about doing the same things faster. It\u2019s about operating at a level that wasn\u2019t previously possible.",
            font_size=15, font_color=BLACK, italic=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# 3. WHERE WE ARE
# ============================================================
slide = content_slide("where we are")

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(4.5), [
    "Phase 1 complete \u2014 AI maintenance assistant live on the shop floor",
    "Dave\u2019s expertise captured: 1,241 documents indexed, 570-vertex knowledge graph",
    "Real machinists asking real questions (thrust bearings, thread peaks, short parts)",
    "System delivers step-by-step, source-cited guidance in seconds",
    "Currently in the Adopt stage \u2014 system works, needs more knowledge and more pilots",
], font_size=20, spacing_pt=14)

add_callout_box(slide, Inches(8.5), Inches(5.5), Inches(3.8), Inches(0.9),
                "Status: WHEELS UP", font_size=20)


# ============================================================
# 4. THE OPPORTUNITY
# ============================================================
slide = content_slide("the opportunity")

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(4.0), [
    "Gent has more Davenport capacity than can currently be staffed",
    "Dave has retired \u2014 his knowledge is in the system, not walking around",
    "The bottleneck isn\u2019t equipment \u2014 it\u2019s skilled people",
    "The AI cockpit lets less experienced operators perform at expert level",
    "Each additional Davenport brought online = $120,000\u2013$168,000/year in new revenue",
], font_size=20, spacing_pt=14)

sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.2), Inches(8), Inches(1.0))
sh.fill.solid(); sh.fill.fore_color.rgb = TABLE_ALT_ROW
sh.line.color.rgb = AIDAN_BLUE; sh.line.width = Pt(2)
tf = sh.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Phase 2 isn\u2019t about protecting what you have. It\u2019s about growing what you can do."
p.font.size = Pt(16); p.font.color.rgb = AIDAN_BLUE; p.font.italic = True
p.font.name = FONT_CONTENT; p.alignment = PP_ALIGN.CENTER


# ============================================================
# 5. THE INDUSTRIAL PILOT
# ============================================================
slide = content_slide("the industrial pilot", "How we enable your team")

add_table(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5), [
    ["", "Traditional", "Industrial Pilot (with Aidan)"],
    ["Leadership", "Manager \u2014 monitors and controls", "Flight Instructor \u2014 coaches your team to fly solo with AI"],
    ["Worker", "Employee \u2014 does tasks, waits for the expert", "Pilot \u2014 navigates independently with AI-powered instruments"],
    ["Environment", "Workplace \u2014 clock in, do the job", "Airframe \u2014 the AI platform that everything runs on"],
    ["Culture", "Silos \u2014 knowledge stuck in individual heads", "Shared Airspace \u2014 everyone\u2019s knowledge available to everyone, 24/7"],
    ["Visibility", "Time card \u2014 hours worked", "Altitude Report \u2014 performance visibility across all cockpits"],
], font_size=14, col_widths=[Inches(2), Inches(4.2), Inches(5.3)])


# ============================================================
# 6. THE FLIGHT PATH
# ============================================================
slide = content_slide("the flight path", "Six stages from ground to air")

tbl = add_table(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5), [
    ["Stage", "Flight Term", "Business Term", "Status"],
    ["0", "PRE-FLIGHT", "Assess", "\u2713 Complete"],
    ["1", "CONFIGURE", "Configure & Deploy", "\u2713 Complete"],
    ["2", "TAXI", "Train & Tune", "\u2713 Complete"],
    ["3", "WHEELS UP", "Adopt", "\u2190 You are here"],
    ["4", "CLIMBING", "Expand", "Next"],
    ["5", "CRUISING", "Compound", "Goal"],
], font_size=16, col_widths=[Inches(1), Inches(2.5), Inches(3.5), Inches(4.5)])

for c in range(4):
    cell = tbl.table.cell(4, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = COLORS["main"]
    cell.text_frame.paragraphs[0].font.bold = True


# ============================================================
# 7. THREE VALUE PILLARS
# ============================================================
slide = content_slide("three value pillars")

for i, (t, d) in enumerate([
    ("scale", "Get more machines running.\n\nA junior hire paired with the AI cockpit can operate a Davenport that would otherwise sit idle."),
    ("capability", "Raise the floor across all operators.\n\nEvery machinist gets expert-level troubleshooting \u2014 consistent quality, faster changeover."),
    ("redundancy", "Remove single points of failure.\n\nNo one person is a bottleneck. Every shift has the same knowledge available."),
]):
    left = Inches(0.8) + i * Inches(4.1)
    add_callout_box(slide, left, Inches(1.8), Inches(3.5), Inches(0.8), t, font_size=22)
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.8), Inches(3.5), Inches(3.2))
    body.fill.solid(); body.fill.fore_color.rgb = TABLE_ALT_ROW
    body.line.color.rgb = COLORS["main"]; body.line.width = Pt(1)
    tf = body.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = d
    p.font.size = Pt(15); p.font.color.rgb = BLACK; p.font.name = FONT_CONTENT
    p.alignment = PP_ALIGN.CENTER


# ============================================================
# 8. WHEELS UP — ADOPT
# ============================================================
slide = content_slide("wheels up \u2014 adopt", "Months 1\u20133: What we\u2019re going to do")

add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(4.0), [
    "\"Ask the Cockpit First\" \u2014 team expectation set by Rich",
    "Close known knowledge gaps (extension pin, oil specs, thrust bearing brands)",
    "Get Dave engaged remotely \u2014 reviewing answers on his phone, flagging what\u2019s wrong",
    "Content sprint: cover top 10 troubleshooting scenarios",
    "Shop floor display board \u2014 weekly visibility",
    "Biweekly coaching calls with Rich",
    "Monthly Altitude Report posted on the shop floor",
], font_size=18, spacing_pt=10)

sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8))
sh.fill.solid(); sh.fill.fore_color.rgb = TABLE_ALT_ROW
sh.line.color.rgb = AIDAN_BLUE; sh.line.width = Pt(1)
tf = sh.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Target: 2+ queries per machinist per week  \u00b7  70%+ top scenarios covered  \u00b7  Feedback flowing from Dave and the team"
p.font.size = Pt(14); p.font.color.rgb = AIDAN_BLUE; p.font.bold = True
p.font.name = FONT_CONTENT; p.alignment = PP_ALIGN.CENTER


# ============================================================
# 9. CLIMBING TO CRUISING
# ============================================================
slide = content_slide("climbing to cruising", "Months 4+: Expand and compound")

add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(3.5), [
    "Add setup procedures to KB \u2014 enables less experienced operators",
    "Altitude Report live \u2014 leaderboard, personal stats, team progress",
    "Growth planning: which jobs can a junior hire handle with cockpit support?",
    "When ready: add a New Leg to the airframe ($10,500)",
    "System is comprehensive \u2014 daily use, value compounding",
], font_size=18, spacing_pt=12)

sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8))
sh.fill.solid(); sh.fill.fore_color.rgb = TABLE_ALT_ROW
sh.line.color.rgb = AIDAN_BLUE; sh.line.width = Pt(1)
tf = sh.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Target: 5+ queries/week  \u00b7  70%+ satisfaction  \u00b7  Setup procedures indexed  \u00b7  Growth path identified"
p.font.size = Pt(14); p.font.color.rgb = AIDAN_BLUE; p.font.bold = True
p.font.name = FONT_CONTENT; p.alignment = PP_ALIGN.CENTER


# ============================================================
# 10. SERVICE TIERS
# ============================================================
slide = content_slide("managed air \u2014 service tiers")

tbl = add_table(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(5.2), [
    ["", "Small \u2014 Maintain Altitude", "Medium \u2014 Climb"],
    ["Monthly Fee", "$800/mo", "$1,600/mo"],
    ["Azure Infrastructure", "\u2713 Included", "\u2713 Included"],
    ["SOC 2 Compliance", "\u2713 Included", "\u2713 Included"],
    ["Monthly Altitude Report", "\u2713", "\u2713"],
    ["Minor Adjustments", "\u2713", "\u2713"],
    ["Proactive Gap Analysis", "\u2014", "\u2713"],
    ["Content Creation", "\u2014", "\u2713"],
    ["Agent Tuning", "\u2014", "\u2713"],
    ["Monthly Coaching Call", "\u2014", "\u2713"],
    ["Best For", "Cruising", "Adopt \u2192 Climbing"],
], font_size=15, col_widths=[Inches(3.5), Inches(3.25), Inches(3.25)])

# Highlight Medium header
tbl.table.cell(0, 2).fill.solid()
tbl.table.cell(0, 2).fill.fore_color.rgb = AIDAN_BLUE


# ============================================================
# 11. EXPAND THE AIRFRAME
# ============================================================
slide = content_slide("expand the airframe")

# Full Throttle card
add_callout_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.8),
                "full throttle \u2014 solution expansion", font_size=18)
add_textbox(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(0.5),
            "$1,500 per story point  \u00b7  2\u20134 points ($3,000\u2013$6,000)",
            font_size=15, font_color=AIDAN_BLUE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(3.2), Inches(5), Inches(2.5), [
    "New knowledge areas or deep content buildout",
    "Search quality overhaul or retrieval tuning",
    "Scoped improvement sprint",
], font_size=15, spacing_pt=8)

# New Leg card
add_callout_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.8),
                "new leg \u2014 new work center", font_size=18)
add_textbox(slide, Inches(7.0), Inches(2.6), Inches(5.5), Inches(0.5),
            "$10,500  \u00b7  7 story points at $1,500/point",
            font_size=15, font_color=AIDAN_BLUE, bold=True)
add_bullet_list(slide, Inches(7.2), Inches(3.2), Inches(5), Inches(2.5), [
    "Full Configure for a new machine type on the airframe",
    "Expert interviews, document ingestion, knowledge graph",
    "Agent configuration, tuning, and shop floor testing",
    "Same monthly Managed AIR covers all work centers",
], font_size=15, spacing_pt=8)

add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
            "More legs on the airframe, more value per dollar.",
            font_size=14, font_color=AIDAN_BLUE, italic=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# 12. MANAGED AIR — YOUR AI OPERATING SYSTEM
# ============================================================
slide = content_slide("managed air", "Your AI operating system")

for i, (comp, desc) in enumerate([
    ("monitoring & health", "Continuous surveillance \u2014\ncatch issues before failures"),
    ("training & enablement", "Your people become confident,\ncapable AI users"),
    ("continuous improvement", "Models evolve, outcomes\naccelerate \u2014 value compounds"),
    ("change management", "Sense and respond as\nthe landscape shifts"),
    ("governance", "Policies current, data quality\nmaintained, risks managed"),
    ("data quality", "Clean, current, trustworthy data \u2014\nAI is only as good as\nwhat you feed it"),
]):
    row, col = i // 3, i % 3
    left = Inches(0.8) + col * Inches(4.1)
    top = Inches(2.2) + row * Inches(2.4)
    add_callout_box(slide, left, top, Inches(3.6), Inches(0.7), comp, font_size=15)
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(0.8), Inches(3.6), Inches(1.2))
    body.fill.solid(); body.fill.fore_color.rgb = TABLE_ALT_ROW
    body.line.color.rgb = COLORS["main"]; body.line.width = Pt(1)
    tf = body.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = desc
    p.font.size = Pt(13); p.font.color.rgb = BLACK; p.font.name = FONT_CONTENT
    p.alignment = PP_ALIGN.CENTER


# ============================================================
# 13. WITHOUT VS. WITH
# ============================================================
slide = content_slide("without vs. with managed air")

RED = RGBColor(0x8B, 0x0C, 0x0C)
GREEN = RGBColor(0x0C, 0x6B, 0x2E)

for label, color, left, items, text_color in [
    ("without managed air", RED, Inches(0.8), [
        "System drifts \u2014 no one maintaining or improving",
        "Team reverts to old habits",
        "Knowledge gaps stay open \u2014 trust erodes",
        "Phase 1 investment slowly goes to waste",
    ], RGBColor(0x55, 0x22, 0x22)),
    ("with managed air", GREEN, Inches(7.0), [
        "Performance tracked and reported monthly",
        "Team coached through adoption",
        "Knowledge base grows \u2014 system gets smarter every month",
        "Phase 1 investment compounds into real growth",
    ], RGBColor(0x0C, 0x4B, 0x2E)),
]:
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(5.5), Inches(0.7))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    tf = sh.text_frame
    p = tf.paragraphs[0]; p.text = label
    p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.bold = True
    p.font.name = FONT_CONTENT; p.alignment = PP_ALIGN.CENTER
    add_bullet_list(slide, left + Inches(0.2), Inches(2.5), Inches(5), Inches(3.5),
                    items, font_size=16, spacing_pt=12, font_color=text_color)


# ============================================================
# 14. THE ASK
# ============================================================
slide = content_slide("the ask")

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(2.5), [
    "Start at Medium tier for 6 months to drive adoption",
    "Drop to Small once the system is self-sustaining",
    "Less than the cost of one extended downtime incident per month",
    "When ready to expand: Full Throttle sprints or New Legs on the airframe",
], font_size=20, spacing_pt=14)

add_table(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(2.2), [
    ["Period", "Tier", "Focus"],
    ["Months 1\u20133", "Medium \u2014 Climb", "Drive daily use, close gaps, coaching cadence"],
    ["Months 4\u20136", "Medium \u2014 Climb", "Add setup procedures, Altitude Report, growth planning"],
    ["Month 7+", "Small \u2014 Maintain Altitude", "Sustain, support new hires, enable expansion"],
], font_size=15, col_widths=[Inches(2), Inches(3), Inches(5)])


# ============================================================
# 15. THE ROI — BUILD UP THE MATH
# ============================================================
slide = content_slide("the return on investment")

# Left side: build up the revenue math
add_textbox(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(0.5),
            "One more Davenport running:", font_size=18, font_color=AIDAN_BLUE, bold=True)

add_table(slide, Inches(0.8), Inches(2.0), Inches(6), Inches(2.8), [
    ["", "Low", "High"],
    ["Billing rate", "$85/hr", "$100/hr"],
    ["\u00d7 Productive hours/day", "6.5 hrs", "6.5 hrs"],
    ["\u00d7 Operating days/month", "20 days", "20 days"],
    ["= Revenue per month", "$11,050", "$13,000"],
    ["= Revenue per year", "$132,600", "$156,000"],
], font_size=14, col_widths=[Inches(3), Inches(1.5), Inches(1.5)])

# Managed AIR cost below
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(6), Inches(0.4),
            "Managed AIR cost (Year 1):", font_size=16, font_color=AIDAN_BLUE, bold=True)
add_textbox(slide, Inches(0.8), Inches(5.4), Inches(6), Inches(0.8),
            "6 mo Medium + 6 mo Small = ~$18,720/year",
            font_size=15, font_color=BLACK)

# Right side: the big number
add_callout_box(slide, Inches(8), Inches(1.8), Inches(4.5), Inches(1.8),
                "6\u20139\u00d7", font_size=60)
add_textbox(slide, Inches(8), Inches(3.8), Inches(4.5), Inches(0.5),
            "return on managed air",
            font_size=16, font_color=AIDAN_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# The division
add_textbox(slide, Inches(8), Inches(4.5), Inches(4.5), Inches(0.4),
            "$132K \u00f7 $18.7K = 7.1\u00d7", font_size=15, font_color=BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(8), Inches(4.9), Inches(4.5), Inches(0.4),
            "$156K \u00f7 $18.7K = 8.3\u00d7", font_size=15, font_color=BLACK, alignment=PP_ALIGN.CENTER)

# Plus other benefits
add_textbox(slide, Inches(8), Inches(5.6), Inches(4.5), Inches(1.0),
            "Plus: 5%+ efficiency across existing machines\nPlus: faster onboarding for new hires\nPlus: Dave\u2019s expertise available 24/7",
            font_size=13, font_color=BLACK, italic=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# 16. CLOSING
# ============================================================
sn += 1
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, WHITE)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(3.0))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = (
    "The question isn\u2019t whether you can afford Phase 2.\n\n"
    "It\u2019s how many more machines could be running if the skills "
    "bottleneck were removed \u2014 and the knowledge to do that is "
    "already in the system."
)
p.font.size = Pt(24); p.font.color.rgb = BLACK; p.font.italic = True
p.font.name = FONT_CONTENT; p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(36)

sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.2), Inches(2.3), Inches(0.025))
sh.fill.solid(); sh.fill.fore_color.rgb = AIDAN_BLUE; sh.line.fill.background()

add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
            "Aidan Systems  \u00b7  MAGNET  \u00b7  March 2026",
            font_size=16, font_color=BLACK, alignment=PP_ALIGN.CENTER)

if os.path.exists(LOGO_AIDAN_ICON):
    slide.shapes.add_picture(LOGO_AIDAN_ICON, Inches(6.2), Inches(6.2), height=Inches(0.8))

add_slide_number(slide, sn + 1)


# --- Save ---
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
