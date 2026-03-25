"""
Aidan Systems — Presentation Format Configuration
Load this in any deck script for consistent branding.

Design: Clean white background, Apple-esque minimal.
Fonts: Bauhaus (lowercase) for titles, Century Gothic for content.
Colors: Aidan website palette.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# BRAND COLORS (from aidansystems.com)
# ============================================================
COLORS = {
    "main":       RGBColor(0xE1, 0xEF, 0xF2),  # light green
    "secondary":  RGBColor(0x95, 0xDA, 0xBB),  # grey green
    "black":      RGBColor(0x19, 0x19, 0x19),  # black (text)
    "aidan_blue": RGBColor(0x05, 0x98, 0xB3),  # aidan blue (primary accent)
    "white":      RGBColor(0xFF, 0xFF, 0xFF),  # white
    "dark":       RGBColor(0x00, 0x00, 0x00),  # really black
}

# Convenience aliases
AIDAN_BLUE = COLORS["aidan_blue"]
BLACK = COLORS["black"]
WHITE = COLORS["white"]
MAIN_BG = COLORS["main"]
SECONDARY = COLORS["secondary"]
DARK = COLORS["dark"]

# Table colors
TABLE_HEADER_BG = AIDAN_BLUE
TABLE_HEADER_TEXT = WHITE
TABLE_ALT_ROW = RGBColor(0xF5, 0xF9, 0xFA)  # very light version of main
TABLE_BORDER = RGBColor(0xD0, 0xD0, 0xD0)

# ============================================================
# FONTS
# ============================================================
FONT_TITLE = "Bauhaus"       # Titles — always lowercase
FONT_CONTENT = "Century Gothic"  # Body text, bullets, tables

# ============================================================
# SLIDE DIMENSIONS (widescreen 16:9)
# ============================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ============================================================
# LOGO PATHS (place logos in this directory)
# ============================================================
LOGO_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "branding")
LOGO_AIDAN = os.path.join(LOGO_DIR, "Aidan-Logo.png")       # Aidan full logo
LOGO_AIDAN_ICON = os.path.join(LOGO_DIR, "Aidan-Icon.png")  # Aidan icon (for slide corners)
LOGO_MAGNET = os.path.join(LOGO_DIR, "logo-magnet.png")     # MAGNET logo (add when available)
LOGO_CUSTOMER = os.path.join(LOGO_DIR, "logo-customer.png") # Customer logo (swap per deck)


# ============================================================
# HELPER FUNCTIONS
# ============================================================

def set_slide_bg(slide, color=WHITE):
    """Set slide background. Default: clean white."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, left=Inches(0.8), top=Inches(0.4),
              width=Inches(11), height=Inches(1.0), font_size=60):
    """Add a Bauhaus lowercase title in aidan blue."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text.lower()  # always lowercase for Bauhaus
    p.font.size = Pt(font_size)
    p.font.color.rgb = AIDAN_BLUE
    p.font.bold = False  # Bauhaus is already bold by design
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.LEFT
    return txBox


def add_subtitle(slide, text, left=Inches(0.8), top=Inches(1.3),
                 width=Inches(11), height=Inches(0.6), font_size=20):
    """Add a Century Gothic subtitle in black."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = BLACK
    p.font.bold = False
    p.font.name = FONT_CONTENT
    p.alignment = PP_ALIGN.LEFT
    return txBox


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                italic=False):
    """Add a Century Gothic text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT_CONTENT
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    font_color=BLACK, bold_items=None, spacing_pt=8):
    """Add a bulleted list in Century Gothic."""
    if bold_items is None:
        bold_items = set()
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = FONT_CONTENT
        p.font.bold = (i in bold_items)
        p.space_after = Pt(spacing_pt)
        p.level = 0
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        for child in list(pPr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag.startswith('bu'):
                pPr.remove(child)
        pPr.append(buChar)
    return txBox


def add_table(slide, left, top, width, height, data, col_widths=None,
              font_size=14):
    """Add a table with aidan blue header row and clean alternating rows."""
    rows = len(data)
    cols = len(data[0]) if data else 0
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(cell_text)
            p.font.size = Pt(font_size)
            p.font.name = FONT_CONTENT
            if row_idx == 0:
                p.font.color.rgb = TABLE_HEADER_TEXT
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            else:
                p.font.color.rgb = BLACK
                cell.fill.solid()
                if row_idx % 2 == 1:
                    cell.fill.fore_color.rgb = WHITE
                else:
                    cell.fill.fore_color.rgb = TABLE_ALT_ROW
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def add_accent_line(slide, left=Inches(0.8), top=Inches(1.15),
                    width=Inches(2)):
    """Add a thin aidan blue accent line under a title."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                    width, Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AIDAN_BLUE
    shape.line.fill.background()
    return shape


def add_slide_number(slide, number):
    """Add slide number in bottom right, light gray."""
    add_textbox(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                str(number), font_size=10,
                font_color=RGBColor(0xBB, 0xBB, 0xBB),
                alignment=PP_ALIGN.RIGHT)


def add_callout_box(slide, left, top, width, height, text,
                    bg_color=AIDAN_BLUE, text_color=WHITE, font_size=18):
    """Add a rounded rectangle callout box."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = FONT_CONTENT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_card(slide, left, top, width, height, title_text, body_text,
             title_size=20, body_size=14):
    """Add a card with aidan blue title header and light body."""
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left, top, width, Pt(40))
    header.fill.solid()
    header.fill.fore_color.rgb = AIDAN_BLUE
    header.line.fill.background()
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(title_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_CONTENT
    p.alignment = PP_ALIGN.CENTER

    # Body
    body_top = top + Pt(44)
    body_height = height - Pt(48)
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, body_top, width, body_height)
    body.fill.solid()
    body.fill.fore_color.rgb = TABLE_ALT_ROW
    body.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    body.line.width = Pt(1)
    tf = body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = body_text
    p.font.size = Pt(body_size)
    p.font.color.rgb = BLACK
    p.font.name = FONT_CONTENT
    p.alignment = PP_ALIGN.CENTER
    return header, body


def add_icon(slide, left=Inches(12.2), top=Inches(0.3), height=Inches(0.4)):
    """Add small Aidan icon to top-right of slide."""
    if os.path.exists(LOGO_AIDAN_ICON):
        slide.shapes.add_picture(LOGO_AIDAN_ICON, left, top, height=height)


def add_logos(slide, aidan=True, magnet=False, customer=False,
             top=Inches(6.8), height=Inches(0.5)):
    """Add logos to bottom of slide if files exist."""
    x = Inches(0.8)
    for show, path in [(aidan, LOGO_AIDAN), (magnet, LOGO_MAGNET),
                        (customer, LOGO_CUSTOMER)]:
        if show and os.path.exists(path):
            slide.shapes.add_picture(path, x, top, height=height)
            x += Inches(2.5)
